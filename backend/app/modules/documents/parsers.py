import io
import logging
import shutil
from typing import List, Optional
from pydantic import BaseModel
import fitz
from pptx import Presentation
from docx import Document
from openpyxl import load_workbook


logger = logging.getLogger(__name__)


ALLOWED_EXTENSIONS = {"pdf", "pptx", "docx", "xlsx"}
OCR_MIN_TEXT_CHARS = 20


MAGIC_BYTES = {
    b"%PDF": "pdf",
    b"PK\x03\x04": "pptx",
}


class PdfPasswordRequired(ValueError):
    """Raised when an encrypted PDF needs a password."""


class PdfPasswordInvalid(ValueError):
    """Raised when a provided PDF password cannot unlock the document."""


class OcrRuntimeUnavailable(RuntimeError):
    """Raised when local OCR dependencies are not available."""


class ParsedSection(BaseModel):
    title: str
    body: str
    position: int


class ParseResult(BaseModel):
    material_id: int
    sections: List[ParsedSection]


def validate_file_type(content: bytes, filename: str) -> Optional[str]:
    """Validate file type using magic number detection."""
    if not filename or "." not in filename:
        return None

    ext = filename.split(".")[-1].lower()
    if ext not in ALLOWED_EXTENSIONS:
        return None

    magic = content[:4]
    if ext == "pdf" and magic.startswith(b"%PDF"):
        return ext
    if ext in {"pptx", "docx", "xlsx"} and magic.startswith(b"PK\x03\x04"):
        return ext
    return None


def parse_pdf_content(content: bytes, password: Optional[str] = None) -> List[ParsedSection]:
    sections = []
    doc = None
    try:
        doc = fitz.open(stream=content, filetype="pdf")
        if doc.needs_pass:
            if not password:
                raise PdfPasswordRequired("PDF requires password")
            if not doc.authenticate(password):
                raise PdfPasswordInvalid("Invalid PDF password")

        logger.info(f"PDF has {len(doc)} pages")

        for i, page in enumerate(doc):
            text = extract_page_text(page)
            if len(text.strip()) < OCR_MIN_TEXT_CHARS:
                logger.info(f"Page {i + 1} has insufficient text; trying OCR")
                ocr_text = extract_page_text_with_ocr(page, i + 1)
                if ocr_text.strip():
                    text = ocr_text

            section = section_from_text(text, len(sections), f"Page {i + 1}")
            if section:
                sections.append(section)

        logger.info(f"Total sections extracted: {len(sections)}")
    except (PdfPasswordRequired, PdfPasswordInvalid, OcrRuntimeUnavailable):
        raise
    except Exception as e:
        logger.error(f"PDF parsing error: {e}")
    finally:
        if doc is not None:
            doc.close()
    return sections


def extract_page_text(page) -> str:
    text = page.get_text("text") or ""
    if text.strip():
        return text

    text_dict = page.get_text("dict")
    texts = []
    for block in text_dict.get("blocks", []):
        for line in block.get("lines", []):
            for span in line.get("spans", []):
                if span.get("text"):
                    texts.append(span["text"])
    return "\n".join(texts)


def extract_page_text_with_ocr(page, page_number: int, lang: str = "eng") -> str:
    """Extract English text from a rendered PDF page using local Tesseract."""
    if shutil.which("tesseract") is None:
        raise OcrRuntimeUnavailable(
            "Tesseract OCR runtime is not available. Install the tesseract binary."
        )

    try:
        import pytesseract
        from PIL import Image
    except ImportError as exc:
        raise OcrRuntimeUnavailable(
            "Tesseract OCR Python dependencies are not available. Install pytesseract and Pillow."
        ) from exc

    mat = fitz.Matrix(2.0, 2.0)
    pix = page.get_pixmap(matrix=mat)
    image = Image.open(io.BytesIO(pix.tobytes("png")))
    text = pytesseract.image_to_string(image, lang=lang)
    logger.info(f"OCR page {page_number}: extracted {len(text)} chars")
    return text


def section_from_text(text: str, position: int, fallback_title: str) -> Optional[ParsedSection]:
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    if not lines:
        return None

    return ParsedSection(
        title=(lines[0] or fallback_title)[:200],
        body="\n".join(lines[1:]),
        position=position,
    )


def parse_pptx_content(content: bytes) -> List[ParsedSection]:
    sections = []
    try:
        with io.BytesIO(content) as f:
            prs = Presentation(f)
            for i, slide in enumerate(prs.slides):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        texts.append(shape.text)
                if texts:
                    title = texts[0][:200] if texts else f"Slide {i + 1}"
                    body = "\n".join(texts[1:]) if len(texts) > 1 else ""
                    sections.append(ParsedSection(
                        title=title,
                        body=body,
                        position=i
                    ))
    except Exception as e:
        logger.error(f"PPTX parsing error: {e}")
    return sections


def parse_docx_content(content: bytes) -> List[ParsedSection]:
    sections = []
    try:
        with io.BytesIO(content) as f:
            doc = Document(f)
            current_title = ""
            current_body = []
            position = 0
            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                if para.style.name.startswith("Heading"):
                    if current_title or current_body:
                        sections.append(ParsedSection(
                            title=current_title or f"Section {position + 1}",
                            body="\n".join(current_body),
                            position=position
                        ))
                        position += 1
                        current_body = []
                    current_title = text[:200]
                else:
                    current_body.append(text)
            if current_title or current_body:
                sections.append(ParsedSection(
                    title=current_title or f"Section {position + 1}",
                    body="\n".join(current_body),
                    position=position
                ))
    except Exception as e:
        logger.error(f"DOCX parsing error: {e}")
    return sections


def parse_xlsx_content(content: bytes) -> List[ParsedSection]:
    sections = []
    try:
        with io.BytesIO(content) as f:
            wb = load_workbook(f, read_only=True)
            for sheet_idx, sheet in enumerate(wb.sheetnames):
                ws = wb[sheet]
                rows = list(ws.iter_rows(values_only=True))
                if not rows:
                    continue
                headers = rows[0] if rows else ()
                data_rows = rows[1:] if len(rows) > 1 else []
                table_data = []
                for row in data_rows:
                    if any(cell is not None for cell in row):
                        try:
                            table_data.append(dict(zip(headers, row)))
                        except Exception:
                            pass
                if table_data:
                    body_lines = [f"Sheet: {sheet}"]
                    for row in table_data[:20]:
                        row_str = " | ".join(f"{k}: {v}" for k, v in row.items() if v)
                        body_lines.append(row_str)
                    sections.append(ParsedSection(
                        title=sheet[:200],
                        body="\n".join(body_lines),
                        position=sheet_idx
                    ))
    except Exception as e:
        logger.error(f"XLSX parsing error: {e}")
    return sections


def parse_material(
    file_type: str, content: bytes, password: Optional[str] = None
) -> Optional[List[ParsedSection]]:
    """Parse material content into sections.
    
    Args:
        file_type: File extension (e.g., 'pdf', 'pptx')
        content: Raw file content in bytes
        
    Returns:
        List of ParsedSection objects, or None if file type not allowed
    """
    handlers = {
        "pdf": parse_pdf_content,
        "pptx": parse_pptx_content,
        "docx": parse_docx_content,
        "xlsx": parse_xlsx_content,
    }
    handler = handlers.get(file_type.lower())
    if handler:
        if file_type.lower() == "pdf":
            return handler(content, password=password)
        return handler(content)
    return None
